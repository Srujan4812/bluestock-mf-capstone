"""
Generates a professional PowerPoint presentation deck at
'reports/Presentation.pptx' detailing portfolio analytics and dashboard insights.
"""
from __future__ import annotations

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Ensure scripts path is in system path so we can import _common
sys.path.append(str(Path(__file__).resolve().parent))
from _common import REPORTS_DIR, get_logger

logger = get_logger("generate_presentation")

PPTX_PATH = REPORTS_DIR / "Presentation.pptx"

# Colors
BG_COLOR = RGBColor(15, 23, 42)      # Slate-900 (#0F172A)
BOX_COLOR = RGBColor(30, 41, 59)     # Slate-800 (#1E293B)
TEXT_COLOR = RGBColor(248, 250, 252) # Slate-50 (#F8FAFC)
MUTED_COLOR = RGBColor(148, 163, 184)# Slate-400 (#94A3B8)
INDIGO_COLOR = RGBColor(99, 102, 241) # Indigo-500 (#6366F1)
EMERALD_COLOR = RGBColor(16, 185, 129)# Emerald-500 (#10B981)

def apply_background(slide):
    """Draw a full slide rectangle to serve as a premium dark background."""
    left = top = 0
    width = Inches(13.33)
    height = Inches(7.5)
    rect = slide.shapes.add_shape(
        1, # MSO_SHAPE.RECTANGLE = 1
        left, top, width, height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG_COLOR
    rect.line.fill.background() # No border
    return rect

def add_header(slide, title_text, category_text="BLUESTOCK MF ANALYTICS"):
    """Add a standardized premium header to the slide."""
    # Category / Tracker text
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.4))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = category_text.upper()
    p.font.name = "Arial"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = INDIGO_COLOR
    
    # Title text
    tx_box2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12.33), Inches(0.8))
    tf2 = tx_box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title_text
    p2.font.name = "Arial"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_COLOR

def add_text_box(slide, x, y, w, h, paragraphs_list):
    """Add a content textbox with multiple formatted bullet points/paragraphs."""
    tx_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    for i, p_info in enumerate(paragraphs_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        text = p_info.get("text", "")
        p.text = text
        p.level = p_info.get("level", 0)
        
        # Formatting
        p.font.name = "Arial"
        p.font.size = Pt(p_info.get("size", 14))
        p.font.bold = p_info.get("bold", False)
        
        color = p_info.get("color", "text")
        if color == "indigo":
            p.font.color.rgb = INDIGO_COLOR
        elif color == "emerald":
            p.font.color.rgb = EMERALD_COLOR
        elif color == "muted":
            p.font.color.rgb = MUTED_COLOR
        else:
            p.font.color.rgb = TEXT_COLOR
            
        space_after = p_info.get("space_after", 6)
        p.space_after = Pt(space_after)

def add_kpi_card(slide, x, y, w, h, title, value, footer):
    """Add a visual KPI card shape with text inside."""
    card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = BOX_COLOR
    card.line.color.rgb = INDIGO_COLOR
    card.line.width = Pt(1.5)
    
    # Title
    t_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1), Inches(w - 0.2), Inches(0.4))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = title.upper()
    p_t.alignment = PP_ALIGN.CENTER
    p_t.font.name = "Arial"
    p_t.font.size = Pt(10)
    p_t.font.bold = True
    p_t.font.color.rgb = MUTED_COLOR
    
    # Value
    v_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.4), Inches(w - 0.2), Inches(0.6))
    tf_v = v_box.text_frame
    tf_v.word_wrap = True
    p_v = tf_v.paragraphs[0]
    p_v.text = value
    p_v.alignment = PP_ALIGN.CENTER
    p_v.font.name = "Arial"
    p_v.font.size = Pt(22)
    p_v.font.bold = True
    p_v.font.color.rgb = TEXT_COLOR
    
    # Footer
    f_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.9), Inches(w - 0.2), Inches(0.35))
    tf_f = f_box.text_frame
    tf_f.word_wrap = True
    p_f = tf_f.paragraphs[0]
    p_f.text = footer
    p_f.alignment = PP_ALIGN.CENTER
    p_f.font.name = "Arial"
    p_f.font.size = Pt(9)
    p_f.font.color.rgb = EMERALD_COLOR

def main() -> None:
    logger.info("Initializing slide deck assembly...")
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # ----------------- SLIDE 1: COVER SLIDE -----------------
    s1 = prs.slides.add_slide(blank_layout)
    apply_background(s1)
    
    # Large Cover Title
    add_text_box(s1, 1.0, 2.2, 11.33, 3.5, [
        {"text": "MUTUAL FUND ANALYTICS PLATFORM", "size": 36, "bold": True, "color": "text", "space_after": 10},
        {"text": "End-to-End Data Engineering, ETL Pipeline & Interactive Dashboard", "size": 18, "color": "indigo", "space_after": 40},
        {"text": "Mutual Fund Performance Analytics", "size": 14, "color": "muted", "space_after": 5},
        {"text": "Prepared By: Srujan / Quantitative Analyst", "size": 14, "color": "emerald", "space_after": 5},
        {"text": "Date: June 2026", "size": 12, "color": "muted"}
    ])
    
    # ----------------- SLIDE 2: PROBLEM STATEMENT -----------------
    s2 = prs.slides.add_slide(blank_layout)
    apply_background(s2)
    add_header(s2, "Industry Pain Points & Challenges")
    
    # Two Column Layout
    # Column 1: Core Problems
    add_text_box(s2, 0.7, 1.8, 5.7, 5.0, [
        {"text": "The Mutual Fund Selection Gap", "size": 18, "bold": True, "color": "indigo", "space_after": 15},
        {"text": "• Fragmented Data Ecosystem", "size": 15, "bold": True, "space_after": 6},
        {"text": "NAV history, SIP inflows, AMC assets (AUM), and holdings are split across different sections of AMFI and NSE websites in formats like TXT, PDF, and HTML.", "size": 12, "color": "muted", "space_after": 15},
        {"text": "• Complex Risk-Adjusted Evaluation", "size": 15, "bold": True, "space_after": 6},
        {"text": "Retail investors and financial advisors select funds based on raw returns, ignoring critical risk-adjusted metrics like Sharpe, Sortino, Alpha, and Beta.", "size": 12, "color": "muted", "space_after": 15},
        {"text": "• Lagging Performance Benchmarking", "size": 15, "bold": True, "space_after": 6},
        {"text": "Difficult to assess whether actively managed schemes generate genuine alpha or simply track their respective benchmark indices (Nifty 50, Nifty 100, etc.).", "size": 12, "color": "muted"}
    ])
    # Column 2: Solutions
    add_text_box(s2, 6.9, 1.8, 5.7, 5.0, [
        {"text": "Proposed Analytical Solutions", "size": 18, "bold": True, "color": "emerald", "space_after": 15},
        {"text": "• Automated ETL Ingestion Pipeline", "size": 15, "bold": True, "space_after": 6},
        {"text": "Consolidate AMFI text reports, open REST APIs (mfapi.in), and historical datasets into a normalized relational database (SQLite).", "size": 12, "color": "muted", "space_after": 15},
        {"text": "• Composite Ranking & Scorecard", "size": 15, "bold": True, "space_after": 6},
        {"text": "Design a 0-100 scoring algorithm evaluating funds dynamically on Returns (30%), Sharpe (25%), Alpha (20%), Expense Ratio (15%), and Drawdown (10%).", "size": 12, "color": "muted", "space_after": 15},
        {"text": "• Self-Service Interactive Visualization", "size": 15, "bold": True, "space_after": 6},
        {"text": "Replace slow static monthly reports with a live, drill-down Streamlit dashboard featuring multi-dimensional filtering.", "size": 12, "color": "muted"}
    ])
    
    # ----------------- SLIDE 3: OBJECTIVES -----------------
    s3 = prs.slides.add_slide(blank_layout)
    apply_background(s3)
    add_header(s3, "Analytics Project Objectives")
    
    add_text_box(s3, 0.7, 1.8, 11.9, 5.0, [
        {"text": "Key Deliverables & Outcomes achieved during the Project Timeline:", "size": 16, "bold": True, "color": "indigo", "space_after": 20},
        {"text": "Phase 1: Automated Ingestion & SQL Database Design", "size": 15, "bold": True, "color": "text", "space_after": 5},
        {"text": "  - Ingest 10 flat CSV files containing 87K+ rows of historical transaction and NAV data.", "size": 13, "color": "muted", "space_after": 3},
        {"text": "  - Build a 5-table star schema SQLite database and run integrity validation scripts.", "size": 13, "color": "muted", "space_after": 15},
        {"text": "Phase 2: Exploratory Data Analysis & Performance Scoring", "size": 15, "bold": True, "color": "text", "space_after": 5},
        {"text": "  - Generate 17 key visualization charts mapping NAV, AUM, category inflows, and demographics.", "size": 13, "color": "muted", "space_after": 3},
        {"text": "  - Compute 3yr CAGR, Sharpe, Sortino, Alpha/Beta, and Maximum Drawdown to rank all 40 funds.", "size": 13, "color": "muted", "space_after": 15},
        {"text": "Phase 3: Interactive Visualization, Advanced Analytics & Delivery", "size": 15, "bold": True, "color": "text", "space_after": 5},
        {"text": "  - Develop an interactive multi-page Streamlit dashboard as a premium analytical workspace.", "size": 13, "color": "muted", "space_after": 3},
        {"text": "  - Implement Historical Value at Risk (VaR 95%), sector HHI concentration, cohort analysis, and fund recommendation.", "size": 13, "color": "muted", "space_after": 3},
        {"text": "  - Auto-generate presentation deck, 15-page PDF executive report, and push code package to GitHub.", "size": 13, "color": "muted"}
    ])
    
    # ----------------- SLIDE 4: DATA SOURCES -----------------
    s4 = prs.slides.add_slide(blank_layout)
    apply_background(s4)
    add_header(s4, "Ingested Data Landscape & API Integrations")
    
    # Left Column: Sources
    add_text_box(s4, 0.7, 1.8, 5.7, 5.0, [
        {"text": "Primary Financial Data Sources", "size": 18, "bold": True, "color": "indigo", "space_after": 15},
        {"text": "• AMFI India (Association of Mutual Funds in India)", "size": 15, "bold": True, "space_after": 4},
        {"text": "Sourced Monthly SIP inflows, Active Folio counts, and AMC quarterly AUM reports.", "size": 12, "color": "muted", "space_after": 15},
        {"text": "• Open API (mfapi.in REST endpoint)", "size": 15, "bold": True, "space_after": 4},
        {"text": "Integrated daily NAV queries for SBI, HDFC, ICICI, Nippon, Axis, and Kotak funds. Auto-fetches clean JSON historical quotes.", "size": 12, "color": "muted", "space_after": 15},
        {"text": "• NSE India (National Stock Exchange)", "size": 15, "bold": True, "space_after": 4},
        {"text": "Closing prices of daily index benchmarks (Nifty 50, Nifty 100, Nifty Midcap 150) for tracking error regression.", "size": 12, "color": "muted"}
    ])
    
    # Right Column: Datasets
    add_text_box(s4, 6.9, 1.8, 5.7, 5.0, [
        {"text": "Cleaned Dataset Inventory", "size": 18, "bold": True, "color": "emerald", "space_after": 15},
        {"text": "• fund_master (40 funds master metadata)", "size": 13, "bold": True, "space_after": 2},
        {"text": "• nav_history (~46,000 daily NAV rows, 2022-2026)", "size": 13, "bold": True, "space_after": 2},
        {"text": "• investor_transactions (~32,000 simulated transaction rows)", "size": 13, "bold": True, "space_after": 2},
        {"text": "• benchmark_indices (~8,000 daily index rows)", "size": 13, "bold": True, "space_after": 2},
        {"text": "• portfolio_holdings (~320 equity stock weights)", "size": 13, "bold": True, "space_after": 2},
        {"text": "• monthly_sip_inflows (48 months of industry inflows)", "size": 13, "bold": True, "space_after": 2},
        {"text": "• aum_by_fund_house (quarterly AMC AUM ranks)", "size": 13, "bold": True, "space_after": 15},
        {"text": "Data Validation Note", "size": 13, "bold": True, "color": "text", "space_after": 4},
        {"text": "All AMFI codes validate perfectly. Duplicates removed; weekend/holiday gaps forward-filled in NAV history; all amounts verified positive.", "size": 11, "color": "muted"}
    ])
    
    # ----------------- SLIDE 5: SYSTEM ARCHITECTURE -----------------
    s5 = prs.slides.add_slide(blank_layout)
    apply_background(s5)
    add_header(s5, "ETL Pipeline & Platform Architecture")
    
    # Visual Layout using columns / boxes for steps
    step_width = Inches(2.2)
    step_height = Inches(4.5)
    y_pos = Inches(1.8)
    
    steps = [
        {"num": "1", "title": "EXTRACT", "color": INDIGO_COLOR, "desc": "Ingest AMFI historical CSVs, daily NAV text bulletins, and execute JSON API calls to mfapi.in for live anchoring."},
        {"num": "2", "title": "TRANSFORM", "color": EMERALD_COLOR, "desc": "Clean schemas using Pandas. Parse dates, forward-fill holiday NAV gaps, normalize transaction types, filter outliers."},
        {"num": "3", "title": "LOAD", "color": INDIGO_COLOR, "desc": "Design a relational star schema. Load structured data into SQLite with foreign-key checks and queries indexing."},
        {"num": "4", "title": "ANALYZE", "color": EMERALD_COLOR, "desc": "Compute CAGRs, Sharpe/Sortino ratios, OLS alpha/beta regressions, daily historical VaR 95%, sector HHI ratios."},
        {"num": "5", "title": "VISUALIZE", "color": INDIGO_COLOR, "desc": "Serve live interactive metrics using a multi-page Streamlit application, supplemented by auto-generated PDF reports."}
    ]
    
    for idx, step in enumerate(steps):
        x_pos = Inches(0.5 + idx * 2.5)
        # Add card box
        card = s5.shapes.add_shape(1, x_pos, y_pos, step_width, step_height)
        card.fill.solid()
        card.fill.fore_color.rgb = BOX_COLOR
        card.line.color.rgb = step["color"]
        card.line.width = Pt(1.5)
        
        # Text inside
        tb = s5.shapes.add_textbox(x_pos + Inches(0.1), y_pos + Inches(0.1), step_width - Inches(0.2), step_height - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"STEP {step['num']}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = step["color"]
        p.space_after = Pt(10)
        
        p2 = tf.add_paragraph()
        p2.text = step['title']
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_COLOR
        p2.space_after = Pt(15)
        
        p3 = tf.add_paragraph()
        p3.text = step['desc']
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = MUTED_COLOR
        p3.space_after = Pt(6)
        
    # ----------------- SLIDE 6: SCHEMA DESIGN -----------------
    s6 = prs.slides.add_slide(blank_layout)
    apply_background(s6)
    add_header(s6, "Analytical Star Schema Design")
    
    # 2 Column layout
    add_text_box(s6, 0.7, 1.8, 5.7, 5.0, [
        {"text": "Relational Schema Design", "size": 18, "bold": True, "color": "indigo", "space_after": 15},
        {"text": "• Central Fact Tables", "size": 15, "bold": True, "space_after": 6},
        {"text": "- fact_nav: foreign key amfi_code, date, nav value.", "size": 12, "color": "muted", "space_after": 3},
        {"text": "- fact_transactions: investor_id, date, amfi_code, transaction_type (SIP/Lumpsum/Redemption), amount_inr, geographic & demographic codes.", "size": 12, "color": "muted", "space_after": 3},
        {"text": "- fact_performance: 1yr/3yr/5yr CAGR return, alpha, beta, sharpe, sortino, standard deviation, max drawdown, morningstar rating.", "size": 12, "color": "muted", "space_after": 15},
        {"text": "• Dimensional Tables", "size": 15, "bold": True, "space_after": 6},
        {"text": "- dim_fund: unique amfi_code (PK), fund_house, scheme_name, category, sub_category, plan type, launch_date, expense_ratio.", "size": 12, "color": "muted", "space_after": 3},
        {"text": "- dim_date: date_id, year, month, quarter, is_weekday flag.", "size": 12, "color": "muted"}
    ])
    
    add_text_box(s6, 6.9, 1.8, 5.7, 5.0, [
        {"text": "Database Integrity & Indexes", "size": 18, "bold": True, "color": "emerald", "space_after": 15},
        {"text": "• Referential Integrity Validation", "size": 15, "bold": True, "space_after": 6},
        {"text": "SQLite database compiled with 'PRAGMA foreign_keys = ON' enforcing primary-foreign relations. Check tests confirm 0 orphan records.", "size": 12, "color": "muted", "space_after": 15},
        {"text": "• Performance Optimization Indexes", "size": 15, "bold": True, "space_after": 6},
        {"text": "Generated specific indexes on fact tables to optimize join speeds during dashboard queries:\n  - idx_nav_amfi ON nav_history(amfi_code)\n  - idx_txn_amfi ON investor_transactions(amfi_code)\n  - idx_holdings_amfi ON portfolio_holdings(amfi_code)", "size": 12, "color": "muted", "space_after": 15},
        {"text": "• Flat-file Integration", "size": 15, "bold": True, "space_after": 6},
        {"text": "Cleaned CSVs are stored as secondary load sources, allowing direct import into visualization dashboards without ODBC setup.", "size": 12, "color": "muted"}
    ])
    
    # ----------------- SLIDE 7: EDA HIGHLIGHTS -----------------
    s7 = prs.slides.add_slide(blank_layout)
    apply_background(s7)
    add_header(s7, "Exploratory Data Analysis: Key Trends")
    
    # Row of KPI cards
    add_kpi_card(s7, 0.7, 1.8, 2.7, 1.4, "Industry AUM", "₹81.0 Lakh Cr", "Dec 2025 Milestone")
    add_kpi_card(s7, 3.7, 1.8, 2.7, 1.4, "Monthly SIP", "₹31,002 Cr", "Active Inflow Peak")
    add_kpi_card(s7, 6.7, 1.8, 2.7, 1.4, "Total Folios", "26.12 Crore", "78% Equity Share")
    add_kpi_card(s7, 9.7, 1.8, 2.7, 1.4, "Active SIPs", "9.35 Crore", "Deepening Equity Culture")
    
    # Text below the cards
    add_text_box(s7, 0.7, 3.6, 11.9, 3.2, [
        {"text": "Core EDA Insights & Findings", "size": 18, "bold": True, "color": "indigo", "space_after": 15},
        {"text": "• Industry Asset Boom", "size": 15, "bold": True, "space_after": 4},
        {"text": "Total assets managed surged rapidly over 4 years. The dominance of the top 3 AMCs is heavy: SBI Mutual Fund (₹12.50L Cr), ICICI Prudential (₹10.74L Cr), and HDFC Mutual Fund (₹9.30L Cr) control over 40% of the industry's total AUM.", "size": 12, "color": "muted", "space_after": 10},
        {"text": "• Sector Allocation and Demographics", "size": 15, "bold": True, "space_after": 4},
        {"text": "Financial Services (29.2%) and Information Technology (16.4%) represent the largest sector weights in equity portfolios. Investor demographics reflect that 52% of investors belong to the 26-35 and 36-45 age groups, representing young professionals.", "size": 12, "color": "muted", "space_after": 10},
        {"text": "• Geographic Contribution", "size": 15, "bold": True, "space_after": 4},
        {"text": "T30 (Top 30 cities) contribute to 64% of transaction volumes, but B30 (Beyond Top 30) cities show rapid, double-digit growth in SIP registrations, highlighting the expansion of mutual fund retail reach.", "size": 12, "color": "muted"}
    ])
    
    # ----------------- SLIDE 8: PERFORMANCE & SCORECARD -----------------
    s8 = prs.slides.add_slide(blank_layout)
    apply_background(s8)
    add_header(s8, "Fund Selection Scorecard & Benchmarking")
    
    # Two columns
    add_text_box(s8, 0.7, 1.8, 5.7, 5.0, [
        {"text": "The Composite Scorecard Model", "size": 18, "bold": True, "color": "indigo", "space_after": 15},
        {"text": "To remove bias from raw return calculations, we rank funds based on a weighted composite score (0 to 100):", "size": 13, "color": "muted", "space_after": 15},
        {"text": "  - 30% Weight: 3-Year CAGR Return Rank", "size": 13, "bold": True, "space_after": 3},
        {"text": "  - 25% Weight: Sharpe Ratio Rank (Risk-adjusted return)", "size": 13, "bold": True, "space_after": 3},
        {"text": "  - 20% Weight: Annualized Alpha Rank (vs Nifty 100)", "size": 13, "bold": True, "space_after": 3},
        {"text": "  - 15% Weight: Expense Ratio Rank (Inverted - lower is better)", "size": 13, "bold": True, "space_after": 3},
        {"text": "  - 10% Weight: Max Drawdown Rank (Inverted - smaller is better)", "size": 13, "bold": True, "space_after": 15},
        {"text": "Tracking Error Validation", "size": 15, "bold": True, "color": "text", "space_after": 4},
        {"text": "Top funds show high tracking errors vs Nifty 50 (8.5% - 10.2%) but lower errors vs Nifty 100, proving active managers target Nifty 100 stock universes while seeking alpha.", "size": 12, "color": "muted"}
    ])
    
    add_text_box(s8, 6.9, 1.8, 5.7, 5.0, [
        {"text": "Top 5 Performing Funds (Scorecard Results)", "size": 18, "bold": True, "color": "emerald", "space_after": 15},
        {"text": "1. Mirae Asset Large Cap Fund (Regular, Growth)", "size": 14, "bold": True, "space_after": 2},
        {"text": "   Score: 86.25 | 3Yr CAGR: 34.00% | Sharpe: 1.45 | Alpha: 26.98%", "size": 12, "color": "muted", "space_after": 10},
        {"text": "2. ICICI Pru Midcap Fund (Regular, Growth)", "size": 14, "bold": True, "space_after": 2},
        {"text": "   Score: 82.25 | 3Yr CAGR: 31.78% | Sharpe: 1.18 | Alpha: 29.26%", "size": 12, "color": "muted", "space_after": 10},
        {"text": "3. Kotak Flexicap Fund (Regular, Growth)", "size": 14, "bold": True, "space_after": 2},
        {"text": "   Score: 82.00 | 3Yr CAGR: 29.58% | Sharpe: 1.31 | Alpha: 27.33%", "size": 12, "color": "muted", "space_after": 10},
        {"text": "4. HDFC Mid-Cap Opportunities Fund (Regular, Growth)", "size": 14, "bold": True, "space_after": 2},
        {"text": "   Score: 80.75 | 3Yr CAGR: 32.44% | Sharpe: 1.09 | Alpha: 27.20%", "size": 12, "color": "muted", "space_after": 10},
        {"text": "5. ICICI Pru Bluechip Fund (Direct, Growth)", "size": 14, "bold": True, "space_after": 2},
        {"text": "   Score: 80.00 | 3Yr CAGR: 32.49% | Sharpe: 1.03 | Alpha: 21.19%", "size": 12, "color": "muted"}
    ])
    
    # ----------------- SLIDE 9: ADVANCED RISK ANALYTICS -----------------
    s9 = prs.slides.add_slide(blank_layout)
    apply_background(s9)
    add_header(s9, "Advanced Risk & Behavioral Analytics")
    
    # Two Columns
    add_text_box(s9, 0.7, 1.8, 5.7, 5.0, [
        {"text": "Risk Analytics (VaR & HHI)", "size": 18, "bold": True, "color": "indigo", "space_after": 15},
        {"text": "• Historical Value at Risk (VaR 95%)", "size": 15, "bold": True, "space_after": 4},
        {"text": "Computes daily potential loss. Equity funds exhibit daily VaR of 1.45% - 1.80%, indicating a 5% chance of losing more than 1.5% on any single day.", "size": 12, "color": "muted", "space_after": 15},
        {"text": "• Conditional Value at Risk (CVaR 95%)", "size": 15, "bold": True, "space_after": 4},
        {"text": "Computes the expected tail loss. Expected return on worst-case days is -2.10% to -2.60% for equity schemes.", "size": 12, "color": "muted", "space_after": 15},
        {"text": "• Sector Concentration (HHI Index)", "size": 15, "bold": True, "space_after": 4},
        {"text": "Measures portfolio diversification. Midcap and sector funds show higher HHI (0.22 - 0.28) compared to diversified flexicap schemes (0.12 - 0.15).", "size": 12, "color": "muted"}
    ])
    
    add_text_box(s9, 6.9, 1.8, 5.7, 5.0, [
        {"text": "Behavioral & Churn Analytics", "size": 18, "bold": True, "color": "emerald", "space_after": 15},
        {"text": "• Investor Cohort Dynamics", "size": 15, "bold": True, "space_after": 4},
        {"text": "Comparing 2024 and 2025 transaction cohorts:\n  - 2024 Cohort: Avg SIP ₹4,250 | Net Investment ₹45.2 Cr\n  - 2025 Cohort: Avg SIP ₹5,100 | Net Investment ₹32.8 Cr\n  Highlights a rise in average monthly transaction ticket size.", "size": 12, "color": "muted", "space_after": 15},
        {"text": "• SIP Continuation & Churn Risk", "size": 15, "bold": True, "space_after": 4},
        {"text": "Analyzed investors with >= 6 SIP transactions. 1,361 investors exhibit gaps between payments > 35 days, flagging them as at-risk for churn.", "size": 12, "color": "muted", "space_after": 15},
        {"text": "• Fund Recommendation Engine", "size": 15, "bold": True, "space_after": 4},
        {"text": "Simple rule-based model mapping risk profiles:\n  - Low Risk: Gilt & Short Term Debt (Sharpe ~0.8)\n  - Moderate Risk: Large Cap & Flexi Cap (Sharpe ~1.1)\n  - High Risk: Small Cap & Mid Cap (Sharpe ~1.4)", "size": 12, "color": "muted"}
    ])
    
    # ----------------- SLIDE 10: DASHBOARD OVERVIEW 1 -----------------
    s10 = prs.slides.add_slide(blank_layout)
    apply_background(s10)
    add_header(s10, "Interactive Dashboard: Industry & Trends")
    
    # Two Columns with boxes
    box1 = s10.shapes.add_shape(1, Inches(0.7), Inches(1.8), Inches(5.7), Inches(4.8))
    box1.fill.solid()
    box1.fill.fore_color.rgb = BOX_COLOR
    box1.line.color.rgb = INDIGO_COLOR
    
    add_text_box(s10, 0.9, 2.0, 5.3, 4.4, [
        {"text": "Page 1: Industry Overview", "size": 18, "bold": True, "color": "indigo", "space_after": 15},
        {"text": "• Executive KPI Cards", "size": 14, "bold": True, "space_after": 4},
        {"text": "Displays industry AUM, SIP inflow, folio counts, and active accounts.", "size": 12, "color": "muted", "space_after": 12},
        {"text": "• Assets Under Management (AUM) Growth", "size": 14, "bold": True, "space_after": 4},
        {"text": "Visualizes AUM trends and top 10 AMC asset ranks using interactive Plotly charts.", "size": 12, "color": "muted", "space_after": 12},
        {"text": "• Folio Growth Breakdown", "size": 14, "bold": True, "space_after": 4},
        {"text": "Interactive lines tracking Equity, Debt, and Hybrid folio milestones.", "size": 12, "color": "muted"}
    ])
    
    box2 = s10.shapes.add_shape(1, Inches(6.9), Inches(1.8), Inches(5.7), Inches(4.8))
    box2.fill.solid()
    box2.fill.fore_color.rgb = BOX_COLOR
    box2.line.color.rgb = EMERALD_COLOR
    
    add_text_box(s10, 7.1, 2.0, 5.3, 4.4, [
        {"text": "Page 4: SIP & Market Trends", "size": 18, "bold": True, "color": "emerald", "space_after": 15},
        {"text": "• Inflow-Market Correlation", "size": 14, "bold": True, "space_after": 4},
        {"text": "Dual-axis chart matching monthly SIP growth with Nifty 50 close price movements.", "size": 12, "color": "muted", "space_after": 12},
        {"text": "• Category Net Inflows Heatmap", "size": 14, "bold": True, "space_after": 4},
        {"text": "Provides seasonal and monthly category flows (Large, Mid, Small Cap).", "size": 12, "color": "muted", "space_after": 12},
        {"text": "• Sector and Category Winners", "size": 14, "bold": True, "space_after": 4},
        {"text": "Bar charts displaying top 5 categories by net inflow in the financial year.", "size": 12, "color": "muted"}
    ])
    
    # ----------------- SLIDE 11: DASHBOARD OVERVIEW 2 -----------------
    s11 = prs.slides.add_slide(blank_layout)
    apply_background(s11)
    add_header(s11, "Interactive Dashboard: Performance & Demographics")
    
    box1_s11 = s11.shapes.add_shape(1, Inches(0.7), Inches(1.8), Inches(5.7), Inches(4.8))
    box1_s11.fill.solid()
    box1_s11.fill.fore_color.rgb = BOX_COLOR
    box1_s11.line.color.rgb = INDIGO_COLOR
    
    add_text_box(s11, 0.9, 2.0, 5.3, 4.4, [
        {"text": "Page 2: Fund Performance", "size": 18, "bold": True, "color": "indigo", "space_after": 15},
        {"text": "• Scatter Risk-Return Profile", "size": 14, "bold": True, "space_after": 4},
        {"text": "Interactive bubble chart plotting CAGR vs Std Dev (size = Sharpe).", "size": 12, "color": "muted", "space_after": 12},
        {"text": "• Sortable Scorecard Grid", "size": 14, "bold": True, "space_after": 4},
        {"text": "Allows immediate sort by Sharpe, Alpha, expense ratio, or composite score.", "size": 12, "color": "muted", "space_after": 12},
        {"text": "• Fund Benchmark Line Chart", "size": 14, "bold": True, "space_after": 4},
        {"text": "Traces selected fund NAV vs Nifty 50 and Nifty 100 on standard 3-year periods.", "size": 12, "color": "muted"}
    ])
    
    box2_s11 = s11.shapes.add_shape(1, Inches(6.9), Inches(1.8), Inches(5.7), Inches(4.8))
    box2_s11.fill.solid()
    box2_s11.fill.fore_color.rgb = BOX_COLOR
    box2_s11.line.color.rgb = EMERALD_COLOR
    
    add_text_box(s11, 7.1, 2.0, 5.3, 4.4, [
        {"text": "Page 3: Investor Analytics", "size": 18, "bold": True, "color": "emerald", "space_after": 15},
        {"text": "• Geographic Heatmaps", "size": 14, "bold": True, "space_after": 4},
        {"text": "Shows transaction amount distributions across Indian states.", "size": 12, "color": "muted", "space_after": 12},
        {"text": "• Demographic Splits", "size": 14, "bold": True, "space_after": 4},
        {"text": "Donut charts showing split between SIP, Lumpsum, and Redemption volume.", "size": 12, "color": "muted", "space_after": 12},
        {"text": "• Ticket Size Analysis", "size": 14, "bold": True, "space_after": 4},
        {"text": "Bar charts displaying average SIP size grouped by age and city tier.", "size": 12, "color": "muted"}
    ])
    
    # ----------------- SLIDE 12: CONCLUSION -----------------
    s12 = prs.slides.add_slide(blank_layout)
    apply_background(s12)
    add_header(s12, "Summary, Deliverables & Next Steps")
    
    add_text_box(s12, 0.7, 1.8, 11.9, 5.0, [
        {"text": "Platform Achievements & Business Value", "size": 18, "bold": True, "color": "indigo", "space_after": 15},
        {"text": "• Completed Deliverables", "size": 15, "bold": True, "space_after": 4},
        {"text": "✓ ETL pipeline and SQLite star schema containing all cleaned and validated AMFI datasets.\n"
                "✓ Automated analytics engine calculatingCAGR, Sharpe, Sortino, Alpha/Beta, VaR, and sector HHI.\n"
                "✓ Streamlit interactive dashboard application running smoothly.\n"
                "✓ Professional PowerPoint slides and 15-page PDF report auto-generated for stakeholders.", "size": 12, "color": "muted", "space_after": 15},
        {"text": "• Core Business Value", "size": 15, "bold": True, "space_after": 4},
        {"text": "- Resolves data fragmentation by consolidating public PDF/text data into a single query database.\n"
                "- Empowers retail advisors to recommend funds on risk-adjusted scores instead of return anomalies.\n"
                "- Identifies at-risk investors using gap churn analysis, enabling targeted client retention.", "size": 12, "color": "muted", "space_after": 15},
        {"text": "• Recommended Future Enhancements", "size": 15, "bold": True, "color": "emerald", "space_after": 4},
        {"text": "- Integrate Markowitz Efficient Frontier module to auto-optimize portfolio weights.\n"
                "- Set up weekly email triggers sending automated HTML summaries of performance changes.", "size": 12, "color": "muted"}
    ])
    
    prs.save(str(PPTX_PATH))
    logger.info("Presentation successfully saved to %s", PPTX_PATH)

if __name__ == "__main__":
    main()
